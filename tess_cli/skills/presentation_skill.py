import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from ..core.logger import setup_logger

logger = setup_logger("PresentationSkill")

class PresentationSkill:
    """
    Skill for creating Stunning Presentations.
    Supports:
    - PPTX (Standard PowerPoint)
    - Marp (Markdown-based slides)
    - Style Selection (Modern, Classic, Tech, Minimal)
    """
    def __init__(self, brain):
        self.brain = brain

    def run(self, action_type, topic, count=5, style="modern", format="pptx", output_name=None):
        """
        Executes the presentation creation.
        """
        try:
            # 1. Generate Content using Brain
            slides_data = self._generate_slides_content(topic, count, style)
            if not slides_data:
                return "Failed to generate slide content."

            # 2. Determine Output Path
            if not output_name:
                output_name = f"{topic.replace(' ', '_')}_Presentation"
            
            # 3. Choose Format
            if format.lower() == "pptx":
                return self.create_pptx(slides_data, output_name, style)
            else:
                return self.create_marp(slides_data, output_name, style)

        except Exception as e:
            logger.error(f"Presentation creation failed: {e}")
            return f"Presentation creation failed: {e}"

    def _generate_slides_content(self, topic, count, style):
        """Uses the Brain to draft content."""
        prompt = (
            f"Draft a {count}-slide presentation about '{topic}' in the style of '{style}'.\n"
            "Return a JSON array of objects, each with 'title' and 'bullets' (list of strings).\n"
            "Keep it professional and visually structured."
        )
        try:
            res = self.brain.think(prompt)
            # Extract JSON
            import re
            match = re.search(r"\[.*\]", res, re.DOTALL)
            if match:
                return json.loads(match.group(0))
        except Exception as e:
            logger.error(f"Brain failed to draft slides: {e}")
        return None

    def create_pptx(self, slides_data, output_name, style="modern"):
        """Creates a .pptx file."""
        prs = Presentation()
        
        # Style Mappings (Basic for now, can be expanded with templates)
        for slide_info in slides_data:
            slide_layout = prs.slide_layouts[1] # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = slide_info.get("title", "Untitled Slide")
            
            tf = content.text_frame
            tf.text = "" # Clear default
            for bullet in slide_info.get("bullets", []):
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0
            
            # Basic Aesthetic Polish based on style
            if style == "tech":
                title.text_frame.paragraphs[0].font.color.rgb = (0, 255, 255) # Cyan
            elif style == "modern":
                title.text_frame.paragraphs[0].font.bold = True

        if not output_name.lower().endswith(".pptx"):
            output_name += ".pptx"
            
        path = os.path.join(os.getcwd(), output_name)
        prs.save(path)
        return f"Successfully created PowerPoint: {path}"

    def create_marp(self, slides_data, output_name, style="gaia"):
        """Creates a Marp-flavored Markdown file."""
        # Marp Header
        marp_content = f"---\nmarp: true\ntheme: {style}\n---\n\n"
        
        for slide in slides_data:
            marp_content += f"# {slide.get('title')}\n\n"
            for bullet in slide.get('bullets', []):
                marp_content += f"- {bullet}\n"
            marp_content += "\n---\n\n"
            
        if not output_name.lower().endswith(".md"):
            output_name += ".md"
            
        path = os.path.join(os.getcwd(), output_name)
        with open(path, "w", encoding="utf-8") as f:
            f.write(marp_content)
        return f"Successfully created Marp Slide Deck: {path}"
